from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

OUT = "/home/obuhow/Рабочий стол/wolf/WOLF_presentation_google-slides.pptx"

# WOLF's quiet register: graphite ink, white surfaces, pale green for progress.
BG = RGBColor(247, 247, 244)
INK = RGBColor(32, 35, 34)
MUTED = RGBColor(102, 108, 104)
RULE = RGBColor(211, 215, 211)
GREEN = RGBColor(190, 221, 195)
GREEN_DARK = RGBColor(61, 107, 72)
GREEN_PALE = RGBColor(232, 242, 232)
WHITE = RGBColor(255, 255, 255)
AMBER = RGBColor(229, 213, 166)
RED_SOFT = RGBColor(235, 214, 211)
FONT = "DejaVu Sans"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_cell_border(cell, color=RULE, width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        tag = qn(edge)
        element = tcPr.find(tag)
        if element is None:
            element = OxmlElement(edge)
            tcPr.append(element)
        element.set('w', width)
        solidFill = element.find(qn('a:solidFill'))
        if solidFill is None:
            solidFill = OxmlElement('a:solidFill')
            element.append(solidFill)
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is None:
            srgb = OxmlElement('a:srgbClr')
            solidFill.append(srgb)
        srgb.set('val', ''.join(f'{v:02X}' for v in color))
        prstDash = element.find(qn('a:prstDash'))
        if prstDash is None:
            prstDash = OxmlElement('a:prstDash')
            element.append(prstDash)
        prstDash.set('val', 'solid')


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, font=FONT, valign=MSO_ANCHOR.TOP,
             margin=0.04, italic=False, line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_rich_text(slide, runs, x, y, w, h, size=18, color=INK, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP, margin=0.04):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for text, bold, col in runs:
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = col
    return box


def rect(slide, x, y, w, h, fill=WHITE, line=RULE, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def line(slide, x1, y1, x2, y2, color=RULE, width=1.0, dash=None):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    if dash:
        ln.line.dash_style = dash
    return ln


def circle(slide, x, y, d, fill=GREEN, line_color=GREEN):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line_color; sh.line.width = Pt(0.75)
    return sh


def title_block(slide, kicker, title, number):
    add_text(slide, kicker.upper(), 0.72, 0.42, 4.5, 0.25, size=9, color=GREEN_DARK, bold=True)
    add_text(slide, title, 0.72, 0.78, 11.7, 0.64, size=27, color=INK, bold=True)
    add_text(slide, f"0{number}", 12.05, 0.43, 0.55, 0.25, size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    line(slide, 0.72, 1.62, 12.62, 1.62, RULE, 0.8)


def footer(slide, text="WOLF · Wolf's Own Life Framework"):
    line(slide, 0.72, 7.08, 12.62, 7.08, RULE, 0.7)
    add_text(slide, text, 0.72, 7.16, 5.8, 0.18, size=8, color=MUTED)


def notes(slide, text):
    try:
        tf = slide.notes_slide.notes_text_frame
        tf.text = text
    except Exception:
        pass


def bullet(slide, text, x, y, w, size=15, color=INK, dot=GREEN_DARK):
    circle(slide, x, y + 0.12, 0.09, dot, dot)
    add_text(slide, text, x + 0.22, y, w - 0.22, 0.35, size=size, color=color)


def arrow(slide, x, y, w, h, fill=GREEN):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill
    return sh

# Slide 1: title / target
slide = prs.slides.add_slide(blank); set_bg(slide)
add_text(slide, "WOLF", 0.72, 0.58, 3.0, 0.65, size=38, color=INK, bold=True)
add_text(slide, "Wolf's Own Life Framework", 0.76, 1.25, 4.8, 0.28, size=12, color=MUTED)
line(slide, 0.76, 1.72, 12.62, 1.72, RULE, 0.8)
add_text(slide, "Для жизни, которая\nне укладывается в расписание", 0.72, 2.18, 7.2, 1.2, size=32, color=INK, bold=True, line_spacing=0.92)
add_text(slide, "Планирование для человека с несколькими проектами,\nволнообразной продуктивностью и автономным графиком.", 0.76, 3.62, 6.4, 0.65, size=17, color=MUTED, line_spacing=1.05)
# right-side life branches
circle(slide, 9.12, 2.58, 1.1, GREEN, GREEN_DARK)
add_text(slide, "Я", 9.12, 2.85, 1.1, 0.32, size=23, color=INK, bold=True, align=PP_ALIGN.CENTER)
for tx, ty, label in [(7.72,1.95,"QA / Java"),(10.48,1.95,"Музыка"),(7.40,4.10,"Здоровье"),(10.70,4.10,"Фриланс"),(9.06,5.10,"Быт")]:
    line(slide, 9.67, 3.12, tx + 0.37, ty + 0.19, GREEN_DARK, 1.1)
    rect(slide, tx, ty, 1.65, 0.46, WHITE, RULE, radius=True)
    add_text(slide, label, tx, ty + 0.10, 1.65, 0.2, size=10, color=INK, align=PP_ALIGN.CENTER)
add_text(slide, "дрейф → разогрев → поток → возвращение", 7.40, 6.07, 4.65, 0.28, size=12, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
footer(slide)
notes(slide, "WOLF создан для людей, которые одновременно развивают несколько направлений и не работают по стабильному внешнему графику. Их реальный путь в работу часто выглядит как дрейф, разогрев, вход в поток, переключение и возвращение. Система не пытается превратить такую жизнь в обычный офисный календарь.")

# Slide 2: audience and problem
slide = prs.slides.add_slide(blank); set_bg(slide); title_block(slide, "01 · Для кого", "Классические трекеры часто усиливают вину", 1)
rect(slide, 0.72, 1.95, 5.78, 4.62, WHITE, RULE, radius=True)
add_text(slide, "Пользователь WOLF", 1.02, 2.22, 3.5, 0.32, size=17, color=INK, bold=True)
bullet(slide, "Живёт сразу в нескольких проектах", 1.02, 2.78, 4.9, size=15)
bullet(slide, "Сам распоряжается своим временем", 1.02, 3.28, 4.9, size=15)
bullet(slide, "Работает волнами, а не равномерно", 1.02, 3.78, 4.9, size=15)
bullet(slide, "Может войти в поток после долгого разогрева", 1.02, 4.28, 4.9, size=15)
bullet(slide, "Часто возвращается к незавершённой работе", 1.02, 4.78, 4.9, size=15)
add_text(slide, "Ему нужен не контроль каждой минуты,\nа опора для возвращения в работу.", 1.02, 5.55, 4.85, 0.62, size=16, color=GREEN_DARK, bold=True)
# right contrast
add_text(slide, "Традиционный сценарий", 7.08, 2.10, 3.7, 0.28, size=13, color=MUTED, bold=True)
for i, (label, col) in enumerate([("Pomodoro", RED_SOFT), ("Жёсткий план дня", RED_SOFT), ("Просрочка = провал", RED_SOFT)]):
    y = 2.62 + i * 0.75
    rect(slide, 7.08, y, 4.66, 0.5, col, col, radius=True)
    add_text(slide, label, 7.32, y + 0.13, 4.15, 0.21, size=14, color=INK)
add_text(slide, "WOLF", 7.08, 5.18, 1.4, 0.28, size=13, color=GREEN_DARK, bold=True)
rect(slide, 7.08, 5.62, 4.66, 0.67, GREEN_PALE, GREEN, radius=True)
add_text(slide, "Недельный объём + реальный темп\n+ возможность вернуться", 7.32, 5.78, 4.15, 0.37, size=15, color=INK, bold=True)
footer(slide)
notes(slide, "Классические инструменты предполагают равномерную продуктивность и внешний ритм. Для пользователя WOLF это часто означает, что отклонение от плана превращается в вину. Поэтому системе нужна другая точка опоры: не контроль каждой минуты, а возможность видеть общий объём и возвращаться к работе.")

# Slide 3: weekly planning
slide = prs.slides.add_slide(blank); set_bg(slide); title_block(slide, "02 · Неделя", "Планируется объём работы, а не каждая минута", 2)
add_text(slide, "Если сложная работа наконец пошла, WOLF не прерывает её\nтолько потому, что закончился формальный слот.", 0.78, 1.93, 8.0, 0.55, size=17, color=MUTED)
# two columns
rect(slide, 0.72, 2.78, 5.55, 3.58, WHITE, RULE, radius=True)
add_text(slide, "План дня", 1.04, 3.08, 1.7, 0.28, size=16, color=MUTED, bold=True)
for i, txt in enumerate(["10:00  задача A", "12:00  задача B", "14:00  задача C"]):
    y = 3.62 + i * 0.65
    rect(slide, 1.04, y, 4.35, 0.42, RED_SOFT, RED_SOFT, radius=True)
    add_text(slide, txt, 1.25, y + 0.10, 3.9, 0.2, size=13, color=INK)
add_text(slide, "Поток приходится\nостанавливать по часам.", 1.04, 5.65, 4.1, 0.43, size=14, color=MUTED)
arrow(slide, 6.43, 4.12, 0.74, 0.42, GREEN)
rect(slide, 7.36, 2.78, 5.2, 3.58, GREEN_PALE, GREEN, radius=True)
add_text(slide, "План недели", 7.68, 3.08, 2.0, 0.28, size=16, color=GREEN_DARK, bold=True)
add_text(slide, "Проект", 7.68, 3.62, 1.7, 0.22, size=11, color=MUTED, bold=True)
add_text(slide, "Объём", 10.42, 3.62, 1.2, 0.22, size=11, color=MUTED, bold=True)
for i, (proj, hrs, done) in enumerate([("QA / Java", "12 ч", 0.60), ("Музыка", "6 ч", 0.36), ("Спорт", "4 ч", 0.84)]):
    y = 4.05 + i * 0.58
    add_text(slide, proj, 7.68, y, 2.2, 0.22, size=13, color=INK)
    add_text(slide, hrs, 10.42, y, 0.7, 0.22, size=13, color=INK, bold=True)
    rect(slide, 11.23, y + 0.05, 0.92, 0.13, WHITE, WHITE, radius=True)
    rect(slide, 11.23, y + 0.05, 0.92 * done, 0.13, GREEN_DARK, GREEN_DARK, radius=True)
add_text(slide, "Распределение внутри недели\nможно менять по факту.", 7.68, 5.72, 4.1, 0.43, size=14, color=GREEN_DARK, bold=True)
footer(slide)
notes(slide, "В WOLF пользователь планирует недельный объём: например, 12 часов на QA, 6 часов на музыку и 4 часа на спорт. Если работа над сложной частью затянулась или, наоборот, пошла быстрее, план можно скорректировать внутри недели. Поток важнее формального окончания слота.")

# Slide 4: forecast
slide = prs.slides.add_slide(blank); set_bg(slide); title_block(slide, "03 · Прогноз", "Как сегодняшние решения меняют будущие сроки", 3)
add_text(slide, "Гантт и кривая нагрузки показывают не только план,\nно и последствия реального темпа.", 0.78, 1.93, 8.5, 0.55, size=17, color=MUTED)
# gantt panel
rect(slide, 0.72, 2.72, 7.62, 3.78, WHITE, RULE, radius=True)
add_text(slide, "Диаграмма Ганта · план / факт / прогноз", 1.02, 2.98, 4.8, 0.25, size=13, color=INK, bold=True)
months = ["Июнь", "Июль", "Авг", "Сент", "Окт"]
for i, m in enumerate(months):
    x = 3.06 + i * 0.92
    add_text(slide, m, x, 3.48, 0.82, 0.18, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    line(slide, x, 3.75, x, 6.02, RULE, 0.6)
rows = [("QA / Java", 0.0, 2.5, GREEN_DARK, "план"), ("Музыка", 0.7, 2.0, GREEN, "факт"), ("Фриланс", 2.1, 2.2, AMBER, "новый проект")]
for i, (name, start, span, col, kind) in enumerate(rows):
    y = 4.02 + i * 0.64
    add_text(slide, name, 1.02, y + 0.09, 1.78, 0.2, size=12, color=INK, bold=(i == 0))
    line(slide, 3.06, y + 0.2, 7.9, y + 0.2, RULE, 0.6)
    rect(slide, 3.06 + start * 0.92, y + 0.06, span * 0.92, 0.28, col, col, radius=True)
    add_text(slide, kind, 3.12 + start * 0.92, y + 0.105, span * 0.92 - 0.1, 0.14, size=8, color=INK, align=PP_ALIGN.CENTER)
# forecast marker
line(slide, 6.55, 3.72, 6.55, 6.05, GREEN_DARK, 1.4)
add_text(slide, "сегодня", 6.22, 6.13, 0.72, 0.18, size=9, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
# load curve panel
rect(slide, 8.65, 2.72, 3.92, 3.78, GREEN_PALE, GREEN, radius=True)
add_text(slide, "Кривая нагрузки", 8.95, 2.98, 2.6, 0.25, size=13, color=GREEN_DARK, bold=True)
# axes
line(slide, 9.10, 5.86, 12.20, 5.86, MUTED, 0.8)
line(slide, 9.10, 5.86, 9.10, 3.64, MUTED, 0.8)
# curve segments
pts = [(9.16,5.54),(9.68,5.18),(10.17,5.32),(10.62,4.40),(11.08,4.78),(11.56,4.12),(12.08,4.28)]
for (x1,y1),(x2,y2) in zip(pts, pts[1:]): line(slide,x1,y1,x2,y2,GREEN_DARK,2.1)
for x,y in pts: circle(slide,x-0.045,y-0.045,0.09,GREEN_DARK,GREEN_DARK)
add_text(slide, "темп", 9.05, 3.41, 0.5, 0.18, size=9, color=MUTED)
add_text(slide, "время", 11.55, 5.98, 0.55, 0.18, size=9, color=MUTED)
add_text(slide, "Новый проект\nсдвигает прогноз", 9.02, 6.07, 2.9, 0.38, size=12, color=INK, bold=True)
footer(slide)
notes(slide, "Пользователь хочет понимать последствия: когда реально закончится QA при текущем темпе, как новый проект повлияет на остальные и где он остановился. Гантт показывает план, факт и прогноз по проектам, а кривая нагрузки помогает увидеть изменение темпа. Это не оценка пользователя, а модель, которую можно пересчитать.")

# Slide 5: ideas
slide = prs.slides.add_slide(blank); set_bg(slide); title_block(slide, "04 · Банк идей", "Идея получает место, но не становится новым долгом", 4)
add_text(slide, "Новые замыслы можно быстро выгрузить из головы,\nне обещая себе немедленно их реализовать.", 0.78, 1.93, 8.2, 0.55, size=17, color=MUTED)
# flow
steps = [("Идея\nпоявилась", 0.9, WHITE), ("Быстрая\nфиксация", 3.25, WHITE), ("Банк\nидей", 5.60, GREEN_PALE), ("Связь с целью\nили сферой", 7.95, WHITE), ("Решение:\nкогда и зачем", 10.30, WHITE)]
for idx, (label, x, fill) in enumerate(steps):
    rect(slide, x, 3.02, 1.72, 1.18, fill, GREEN if fill != WHITE else RULE, radius=True)
    add_text(slide, label, x + 0.12, 3.34, 1.48, 0.48, size=13, color=INK, bold=(idx == 2), align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if idx < len(steps)-1: arrow(slide, x + 1.82, 3.40, 0.40, 0.32, GREEN)
# card
rect(slide, 0.72, 4.82, 5.75, 1.48, WHITE, RULE, radius=True)
add_text(slide, "Пример записи", 1.02, 5.08, 2.0, 0.22, size=11, color=MUTED, bold=True)
add_text(slide, "Автоматизировать личный финансовый учёт", 1.02, 5.40, 4.9, 0.24, size=15, color=INK, bold=True)
add_text(slide, "Сфера: бизнес / быт    ·    Статус: на рассмотрении", 1.02, 5.79, 5.0, 0.2, size=11, color=MUTED)
# right message
rect(slide, 7.00, 4.82, 5.56, 1.48, GREEN_PALE, GREEN, radius=True)
add_text(slide, "Фиксация ≠ обязательство", 7.34, 5.13, 4.8, 0.27, size=17, color=GREEN_DARK, bold=True)
add_text(slide, "Банк идей снижает шум, но не добавляет давления.", 7.34, 5.58, 4.75, 0.28, size=14, color=INK)
footer(slide)
notes(slide, "У пользователя много идей: для бизнеса, творчества, быта и автоматизации. Если не фиксировать их, они продолжают возвращаться и конкурировать с текущей работой. Банк идей отделяет фиксацию от обязательства: идею можно сохранить, связать с целью, а решение о реализации принять позже.")

# Slide 6: meaning layer / close
slide = prs.slides.add_slide(blank); set_bg(slide); title_block(slide, "05 · Слой смыслов", "Связать ежедневную работу с тем, ради чего она нужна", 5)
add_text(slide, "WOLF показывает не только дела и сроки,\nно и место проектов в общей жизни.", 0.78, 1.93, 8.0, 0.55, size=17, color=MUTED)
# network
circle(slide, 5.88, 3.72, 1.12, GREEN, GREEN_DARK)
add_text(slide, "Цели", 5.88, 4.05, 1.12, 0.24, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
items = [("Сферы\nжизни", 1.02, 2.70), ("Проекты", 3.08, 1.98), ("Синергия", 8.35, 1.98), ("LLM Wiki", 10.42, 2.70), ("Банк\nидей", 1.02, 5.12), ("Гантт\nи прогноз", 10.42, 5.12)]
for label, x, y in items:
    rect(slide, x, y, 1.75, 0.72, WHITE, RULE, radius=True)
    add_text(slide, label, x + 0.10, y + 0.17, 1.55, 0.35, size=12, color=INK, bold=(label in ["Проекты", "Гантт\nи прогноз"]), align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    line(slide, 6.44, 4.28, x + 0.87, y + 0.36, GREEN_DARK, 0.9)
# close statement
rect(slide, 3.18, 6.08, 6.96, 0.58, GREEN_PALE, GREEN, radius=True)
add_text(slide, "Не сделать человека идеальной машиной, а помочь ему возвращаться и двигаться в своём реальном темпе.", 3.34, 6.23, 6.64, 0.25, size=13, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
footer(slide)
notes(slide, "Слой смыслов объединяет цели, сферы жизни, проекты, синергию, банк идей, LLM Wiki и честный прогноз. Он отвечает на вопрос: зачем этот проект занимает место в моей жизни и что изменится, если я добавлю ещё один. Итоговая идея WOLF — не сделать человека идеальной машиной продуктивности, а помочь ему возвращаться и двигаться в собственном реальном темпе.")

prs.save(OUT)
print(OUT)
print(f"slides={len(prs.slides)}")
